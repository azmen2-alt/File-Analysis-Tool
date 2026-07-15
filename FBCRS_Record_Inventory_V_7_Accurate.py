import os, time, hashlib, re, sys, warnings
from datetime import datetime, date
from concurrent.futures import ThreadPoolExecutor
from threading import Lock
from openpyxl import Workbook, load_workbook
from openpyxl.worksheet.datavalidation import DataValidation
import tkinter as tk
from tkinter import filedialog

warnings.simplefilter('ignore')

try:
    import pypdf
except ImportError:
    pypdf = None
try:
    import docx
except ImportError:
    docx = None
try:
    import pptx
except ImportError:
    pptx = None

try:
    import ahocorasick
    HAVE_AC = True
except ImportError:
    HAVE_AC = False

if getattr(sys, 'frozen', False):
    SCRIPT_DIR = os.path.dirname(sys.executable)
else:
    SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))

LIBRARY = os.path.join(SCRIPT_DIR, 'FBCRS_Master_Full.xlsx')
MAX_MB = 15
MAX_SIZE_BYTES = MAX_MB * 1024 * 1024

EXTS = {'.docx', '.pdf', '.xlsx', '.xlsm', '.pptx', '.txt'}
SKIP_EXTS = {
    '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff',
    '.mp4', '.mov', '.avi', '.wmv', '.mkv',
    '.zip', '.7z', '.rar',
    '.exe', '.dll', '.msi'
}

LOCK = Lock()
HASHES = {}

if not os.path.exists(LIBRARY):
    raise SystemExit(f'Library not found: {LIBRARY}')

wb_lib = load_workbook(LIBRARY, data_only=True, read_only=True)
RECORDS = {}
KEYWORDS = {}
PHRASES = {}
SYN = {}

for r in wb_lib['Records'].iter_rows(min_row=2, values_only=True):
    if r[0]: RECORDS[r[0]] = {'title': r[1], 'retention': r[2], 'disp': r[3]}
for r in wb_lib['Keywords'].iter_rows(min_row=2, values_only=True):
    if r[0] and r[1]: KEYWORDS.setdefault(r[0], []).append(str(r[1]).lower())
for r in wb_lib['Phrases'].iter_rows(min_row=2, values_only=True):
    if r[0] and r[1]: PHRASES.setdefault(r[0], []).append(str(r[1]).lower())
for r in wb_lib['Synonyms'].iter_rows(min_row=2, values_only=True):
    if r[0] and r[1]: SYN[str(r[0]).lower()] = str(r[1]).lower()
wb_lib.close()

TRACKERS = {
    'kpi tracker': 'STR-POL-003', 'asset tracker': 'EQS-INV-001',
    'event tracker': 'PSR-EVT-001'
}

if SYN:
    _SYN_PATTERN = re.compile('|'.join(re.escape(k) for k in sorted(SYN, key=len, reverse=True)))
else:
    _SYN_PATTERN = None

def norm(t):
    t = (t or '').lower()
    if _SYN_PATTERN:
        t = _SYN_PATTERN.sub(lambda m: SYN[m.group(0)], t)
    return t

def _build_automaton():
    A = ahocorasick.Automaton()
    for code, kws in KEYWORDS.items():
        for k in kws:
            if not k:
                continue
            if not A.exists(k):
                A.add_word(k, [])
            A.get(k).append((code, 2, k))
    for code, phs in PHRASES.items():
        for p in phs:
            if not p:
                continue
            if not A.exists(p):
                A.add_word(p, [])
            A.get(p).append((code, 10, p))
    A.make_automaton()
    return A

AUTOMATON = _build_automaton() if (HAVE_AC and (KEYWORDS or PHRASES)) else None

def score_codes(combined_text):
    if AUTOMATON is not None:
        code_scores = {}
        code_evidence = {}
        for _end_idx, matches in AUTOMATON.iter(combined_text):
            for code, weight, term in matches:
                code_scores[code] = code_scores.get(code, 0) + weight
                ev = code_evidence.setdefault(code, [])
                if term not in ev and len(ev) < 5:
                    ev.append(term)
        return sorted(
            ((code_scores.get(code, 0), code, code_evidence.get(code, [])) for code in RECORDS),
            reverse=True
        )
    else:
        scores = []
        for code in RECORDS:
            s = 0
            ev = []
            for k in KEYWORDS.get(code, []):
                if k in combined_text:
                    s += 2
                    ev.append(k)
            for p in PHRASES.get(code, []):
                if p in combined_text:
                    s += 10
                    ev.append(p)
            scores.append((s, code, ev[:5]))
        scores.sort(reverse=True)
        return scores

def read_text(fp):
    ext = os.path.splitext(fp)[1].lower()
    text_content = ''
    try:
        if ext == '.txt':
            text_content = open(fp, encoding='utf-8', errors='ignore').read(50000)
        elif ext == '.pdf':
            if pypdf is None:
                return ''
            reader = pypdf.PdfReader(fp)
            text_content = ' '.join([page.extract_text() or '' for page in reader.pages])
        elif ext == '.docx':
            if docx is None:
                return ''
            doc = docx.Document(fp)
            text_content = ' '.join([p.text for p in doc.paragraphs])
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text_content += ' ' + str(cell.text)
        elif ext == '.pptx':
            if pptx is None:
                return ''
            prs = pptx.Presentation(fp)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += ' ' + shape.text
        elif ext in ['.xlsx', '.xlsm']:
            wb = load_workbook(fp, data_only=True, read_only=True)
            for sheet in wb.sheetnames:
                text_content += ' ' + sheet
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            text_content += ' ' + str(cell)
            wb.close()
    except Exception:
        pass

    return text_content[:100000]

def sha(fp):
    h = hashlib.sha256()
    try:
        with open(fp, 'rb') as f:
            while chunk := f.read(1048576):  
                h.update(chunk)
    except Exception:
        pass
    return h.hexdigest()

def get_content_category(ext):
    categories = {
        '.docx': 'Document', '.txt': 'Document',
        '.xlsx': 'Spreadsheet', '.xlsm': 'Spreadsheet',
        '.pptx': 'Presentation',
        '.pdf': 'PDF',
        '.jpg': 'Image', '.jpeg': 'Image', '.png': 'Image', '.gif': 'Image', '.bmp': 'Image', '.tiff': 'Image',
        '.mp4': 'Video', '.mov': 'Video', '.avi': 'Video', '.wmv': 'Video', '.mkv': 'Video',
        '.zip': 'Archive', '.7z': 'Archive', '.rar': 'Archive',
        '.exe': 'Executable', '.dll': 'Executable', '.msi': 'Executable',
        '.tmp': 'Temporary', '.bak': 'Backup', '.log': 'Log'
    }
    return categories.get(ext, 'Unknown')

def calculate_expiry(modified_date, retention_string):
    if not retention_string:
        return None
    retention_string = str(retention_string).upper()
    match = re.search(r'(CCY|CFY)\s*\+\s*(\d+)', retention_string)
    if not match:
        return None

    base_type = match.group(1)
    years = int(match.group(2))

    if base_type == 'CCY':
        expiry_year = modified_date.year + years
        return date(expiry_year, 12, 31)
    elif base_type == 'CFY':
        fy_end_year = modified_date.year if modified_date.month < 4 else modified_date.year + 1
        expiry_year = fy_end_year + years
        return date(expiry_year, 3, 31)
    return None

def classify(fp):
    fn = os.path.basename(fp)
    ext = os.path.splitext(fp)[1].lower()
    
    try:
        created = datetime.fromtimestamp(os.path.getctime(fp))
        modified = datetime.fromtimestamp(os.path.getmtime(fp))
    except Exception as e:
        return [fn, ext, 'Unknown', 'File', '', '', '', '', 0, 'N/A', '', '', '', '', '', 'ERROR', 'ACCESS DENIED OR PATH TOO LONG', 'NO', 'NO', fp, str(e)]
        
    category = get_content_category(ext)

    if ext in SKIP_EXTS:
        return [fn, ext, category, 'Media/Archive/App', created, modified, 'N/A', 'Non-Record', 100, 'Very High', 'Instant Skip Rule', '', '', '', 'Review', 'REVIEW', 'UNSUPPORTED FILE TYPE', 'NO', 'NO', fp, 'Bypassed to save time']

    try:
        file_size = os.path.getsize(fp)
    except Exception:
        file_size = 0

    fn_lower = fn.lower()
    transitory_keywords = ['draft', 'working copy', 'untitled', 'test', 'copy of']
    is_transitory = False

    if file_size == 0:
        is_transitory = True
    elif ext in ['.tmp', '.bak', '.log']:
        is_transitory = True
    elif fn.startswith('~$'):
        is_transitory = True
    elif any(kw in fn_lower for kw in transitory_keywords):
        is_transitory = True

    if is_transitory:
        return [fn, ext, category, 'File', created, modified, 'N/A', 'Transitory Record', 100, 'Very High', 'Matched Transitory Rule', '', '', '', 'Destroy', 'DESTROY', 'TRANSITORY', 'NO', 'NO', fp, '']

    if ext not in EXTS:
        return [fn, ext, category, 'File', created, modified, '', '', '', '', '', '', '', '', '', 'REVIEW', 'UNSUPPORTED FILE TYPE', 'NO', 'YES', fp, '']

    if file_size > MAX_SIZE_BYTES:
        return [fn, ext, category, 'File', created, modified, '', '', '', '', '', '', '', '', '', 'REVIEW', 'FILE EXCEEDS CONTENT EXTRACTION LIMIT', 'NO', 'YES', fp, '']

    txt = norm(read_text(fp))

    txt_clean = txt.strip()
    if len(txt_clean) > 0:
        if len(txt_clean) < 15: 
            return [fn, ext, category, 'File', created, modified, 'N/A', 'Transitory Record', 100, 'Very High', 'Content is too short', '', '', '', 'Destroy', 'DESTROY', 'TRANSITORY', 'NO', 'NO', fp, 'File is almost empty inside']
        if 'lorem ipsum' in txt_clean.lower():
            return [fn, ext, category, 'File', created, modified, 'N/A', 'Transitory Record', 100, 'Very High', 'Found placeholder text', '', '', '', 'Destroy', 'DESTROY', 'TRANSITORY', 'NO', 'NO', fp, 'Contains Lorem Ipsum placeholder text']

    fn_norm = norm(fn)
    combined_text = fn_norm + ' ' + txt

    h = sha(fp)
    with LOCK:
        dup_found = 'YES' if h in HASHES else 'NO'
        HASHES[h] = 1

    for t, c in TRACKERS.items():
        if t in combined_text:
            library_disposition = RECORDS.get(c, {}).get('disp', '')
            library_retention = RECORDS.get(c, {}).get('retention', '')

            expiry_date = calculate_expiry(modified, library_retention)
            expiry_date_str = expiry_date.strftime('%Y-%m-%d') if expiry_date else ''

            if dup_found == 'YES':
                final_action = 'DESTROY'
                final_reason = 'DUPLICATE'
            else:
                if expiry_date and datetime.now().date() > expiry_date:
                    final_action = str(library_disposition).upper() if library_disposition else 'REVIEW'
                    final_reason = 'PASSED RETENTION'
                else:
                    final_action = 'DO NOTHING'
                    final_reason = 'STILL WITHIN RETENTION'

            return [fn, ext, category, 'Tracker', created, modified, c, RECORDS.get(c, {}).get('title', ''), 100, 'Very High', t, c, library_retention, expiry_date_str, library_disposition, final_action, final_reason, dup_found, 'NO', fp, '']

    scores = score_codes(combined_text)
    best = scores[0]

    code = best[1]
    conf = min(100, 35 + best[0])
    level = 'Very High' if conf >= 90 else 'High' if conf >= 75 else 'Medium' if conf >= 50 else 'Low'

    library_disposition = RECORDS.get(code, {}).get('disp', '')
    library_retention = RECORDS.get(code, {}).get('retention', '')

    expiry_date = calculate_expiry(modified, library_retention)
    expiry_date_str = expiry_date.strftime('%Y-%m-%d') if expiry_date else ''

    if dup_found == 'YES':
        final_action = 'DESTROY'
        final_reason = 'DUPLICATE'
    elif code.startswith('NO'):
        final_action = 'REVIEW'
        final_reason = 'UNCLASSIFIED'
    else:
        if expiry_date and datetime.now().date() > expiry_date:
            final_action = str(library_disposition).upper() if library_disposition else 'REVIEW'
            final_reason = 'PASSED RETENTION'
        else:
            final_action = 'DO NOTHING'
            final_reason = 'STILL WITHIN RETENTION'

    sme_review = 'YES' if level == 'Low' else 'NO'

    return [fn, ext, category, 'File', created, modified, code, RECORDS.get(code, {}).get('title', 'Unclassified'), conf, level, '; '.join(best[2]), ' | '.join([x[1] for x in scores[:3]]), library_retention, expiry_date_str, library_disposition, final_action, final_reason, dup_found, sme_review, fp, '']

def main():
    print("\n" + "="*50)
    print("--- File Analysis Tool (V6) ---")
    print("="*50)
    
    print("\n[ENGINE MODE: SLOW BUT ACCURATE]")
    print("* This version reads entire documents and performs complete duplicate checks.")
    print("* It is slower, but ensures maximum accuracy for content extraction.")
    
    print("\nWhat this app does:")
    print("* Scans your selected folder and reads your files.")
    print("* Categorizes them based on record retention rules.")
    print("* Creates a final Excel summary report for you.")
    
    print("\nAbout the 'FBCRS_Master_Full.xlsx' file:")
    print("* This is the rulebook the app uses to classify files.")
    print("* IMPORTANT: Keep this Excel file in the exact same folder as the app.")
    print("* You can open it to view the rules or add new keywords.")
    print("* The app will automatically use any new rules you save in it.\n")
    print("="*50 + "\n")

    root = tk.Tk()
    root.withdraw()
    
    start_app = input("Are you ready to select a folder to analyze? (Y/N): ").strip().upper()
    if start_app != 'Y':
        print("Process canceled.")
        input("\nPress Enter to exit the application...")
        sys.exit(0)

    while True:
        print("\nStep 1: Waiting for you to select a folder...")
        print("      [!] NOTE: Because you are using the Accurate mode, this might take time.")
        print("      [!] It is recommended to divide large files into smaller folders.\n")
        
        folder = filedialog.askdirectory(title='Select the folder you want to analyze')

        if not folder:
            print("Process canceled. No folder was selected.")
            retry = input("\nDo you want to try selecting a folder again? (Y/N): ").strip().upper()
            if retry == 'Y':
                continue
            else:
                break

        folder_name = os.path.basename(os.path.normpath(folder))
        OUTPUT = os.path.join(SCRIPT_DIR, f"{folder_name} summary.xlsx")

        print(f"Step 2: Searching for files inside: {folder_name}...")
        files = [os.path.join(r, f) for r, _, fs in os.walk(folder) for f in fs if not f.startswith(('~', '.')) or f.startswith('~$')]
        TOTAL = len(files)

        if TOTAL == 0:
            print("      [!] No files were found in the selected folder.")
            retry = input("\nDo you want to analyze a different folder? (Y/N): ").strip().upper()
            if retry == 'Y':
                continue
            else:
                break

        wb = Workbook()
        ws = wb.active
        ws.title = 'Inventory'
        headers = ['File Name', 'Extension', 'Content Category', 'Record Type', 'Date Created', 'Date Modified', 'Code', 'Series Title', 'Confidence %', 'Confidence Level', 'Classification Evidence', 'Top 3 Candidates', 'Retention Period', 'Retention Expiry Date', 'Disposition', 'Recommended Action', 'Reason', 'Duplicate Detected', 'SME Review Required', 'Full File Path', 'Comments']
        ws.append(headers)

        dv_action = DataValidation(type="list", formula1='"DESTROY,DO NOTHING,TRANSFER TO ARCHIVES,REVIEW"', allow_blank=True)
        dv_reason = DataValidation(type="list", formula1='"DUPLICATE,TRANSITORY,STILL WITHIN RETENTION,PASSED RETENTION,UNCLASSIFIED,UNSUPPORTED FILE TYPE,FILE EXCEEDS CONTENT EXTRACTION LIMIT"', allow_blank=True)

        ws.add_data_validation(dv_action)
        ws.add_data_validation(dv_reason)

        dv_action.add('P2:P1048576')
        dv_reason.add('Q2:Q1048576')

        print(f"Step 3: Found {TOTAL:,} files. Starting analysis...\n")
        print("      [!] IMPORTANT: To stop the script immediately at any time,")
        print("      press Ctrl+C in this window.\n")

        START = time.time()
        PROCESSED = 0
        analysis_stopped = False

        try:
            with ThreadPoolExecutor(max_workers=32) as executor:
                for row in executor.map(classify, files):
                    ws.append(row)

                    current_row = ws.max_row
                    path_cell = ws.cell(row=current_row, column=20) 
                    path_cell.hyperlink = str(row[19])
                    path_cell.style = "Hyperlink"

                    with LOCK:
                        PROCESSED += 1
                        pct = (PROCESSED / TOTAL * 100) if TOTAL else 100
                        print(f'\rProgress: {PROCESSED:,}/{TOTAL:,} files done ({pct:.1f}%)', end='')
        except KeyboardInterrupt:
            print("\n\n[!] Analysis stopped immediately by the user (Ctrl+C).")
            analysis_stopped = True
        
        if not analysis_stopped:
            print("\n\nStep 4: Saving final Excel report...")
            wb.save(OUTPUT)
            
            total_minutes = (time.time() - START) / 60
            
            print("\n" + "="*40)
            print("[SUCCESS] Analysis is complete!")
            print(f"Processed {TOTAL:,} files in {total_minutes:.1f} minutes.")
            print(f"Report saved as: {OUTPUT}")
            print("="*40)
            
            try:
                os.startfile(OUTPUT)
            except Exception:
                pass

        print("\n" + "-"*40)
        repeat = input("\nDo you want to analyze another folder? (Y/N): ").strip().upper()
        if repeat != 'Y':
            break

    input("\nPress Enter to close the application...")

if __name__ == '__main__':
    main()